import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)


class FileParser:
    """Handles parsing of different file formats to extract text content."""
    
    def parse_file(self, file_path: str, file_extension: str) -> str:
        """
        Parse a file and extract text content based on file extension.
        
        Args:
            file_path: Path to the file
            file_extension: File extension (e.g., '.pdf', '.docx')
        
        Returns:
            Extracted text content
        """
        file_extension = file_extension.lower()
        
        if file_extension == '.pdf':
            return self._parse_pdf(file_path)
        elif file_extension in ['.doc', '.docx']:
            return self._parse_doc(file_path)
        elif file_extension in ['.ppt', '.pptx']:
            return self._parse_ppt(file_path)
        else:
            raise ValueError(f"Unsupported file extension: {file_extension}")
    
    def _parse_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            import PyPDF2
            
            text_content = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                
                logger.info(f"Processing PDF with {num_pages} pages")
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
            
            return "\n\n".join(text_content)
        
        except ImportError:
            raise ImportError("PyPDF2 is required for PDF parsing. Install it with: pip install PyPDF2")
        except Exception as e:
            logger.error(f"Error parsing PDF: {str(e)}")
            raise
    
    def _parse_doc(self, file_path: str) -> str:
        """Extract text from DOC/DOCX file."""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n\n".join(text_content)
        
        except ImportError:
            raise ImportError("python-docx is required for DOC/DOCX parsing. Install it with: pip install python-docx")
        except Exception as e:
            logger.error(f"Error parsing DOC/DOCX: {str(e)}")
            raise
    
    def _parse_ppt(self, file_path: str) -> str:
        """Extract text from PPT/PPTX file."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            logger.info(f"Processing PowerPoint with {len(prs.slides)} slides")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
        
        except ImportError:
            raise ImportError("python-pptx is required for PPT/PPTX parsing. Install it with: pip install python-pptx")
        except Exception as e:
            logger.error(f"Error parsing PPT/PPTX: {str(e)}")
            raise

